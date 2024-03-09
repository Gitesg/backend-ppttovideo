from pptx import Presentation

def extract_slide_titles(pptx_file_path):
    # Create a Presentation object
    prs = Presentation(pptx_file_path)
    
    # Dictionary to store slide titles by index
    slide_titles = {}
    i = 0
    
    # Iterate over all slides
    for slide in prs.slides:
        # Check if the slide has a title
        if slide.shapes.title is not None:
            # Get the slide title
            title = slide.shapes.title.text

            # Print the slide title
            print(title)
            
            # Store the slide title in the dictionary
            slide_titles[i] = title
        else:
            # Skip this slide as it does not have a title
            continue
        
        i += 1
    
    return slide_titles

# # Example usage:
# pptx_file_path = 'uploads/W1L2_Switching_Algebra.pptx'
# slide_titles = extract_slide_titles(pptx_file_path)
# print(slide_titles)