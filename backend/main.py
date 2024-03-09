# from flask import Flask, request, jsonify
# import asyncio

# app = Flask(__name__)


# # Define CORS headers middleware
# @app.after_request
# async def add_cors_headers(response):
#     # Allow requests from all origins
#     response.headers.add('Access-Control-Allow-Origin', '*')
#     # Allow specific HTTP methods
#     response.headers.add('Access-Control-Allow-Methods', 'POST')
#     # Allow specific headers
#     response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
#     return response


# # Endpoint for file upload
# @app.route('/upload', methods=['POST'])
# async def upload_file():
#     if 'file' not in request.files:
#         return jsonify({'error': 'No file part'}), 400

#     file = request.files['file']
#     if file.filename == '':
#         return jsonify({'error': 'No selected file'}), 400

#     # Handle file upload asynchronously here
#     # Example: You can use asyncio to handle file upload asynchronously

#     if file and file.filename.endswith(('ppt', 'pptx')):
#         # Perform file upload operations asynchronously here
#         # For demonstration, let's just return a success message
#         return jsonify({'message': 'File uploaded successfully', 'filename': file.filename}), 200
#     else:
#         return jsonify({'error': 'Please upload a PowerPoint file (PPT or PPTX)'}),400
    
# if __name__== '__main__':
#     app.run(debug=True,port=8000)

#-----------------------------------------------------------------------------------------------------------------


from flask import Flask, request, jsonify
import os
import convertapi
from werkzeug.utils import secure_filename
import uuid
from pptx import Presentation
from extracttitles import extract_slide_titles
from text_from_title import generate_explanations
from file_paths import sort_files_and_get_paths
from gemini import generate_image_explanations
from pptx import Presentation
import google.generativeai as genai
from voice import generate_and_save_audio_files
GOOGLE_API_KEY=os.environ.get('gemini')
genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel('gemini-pro-vision')
# import inx

app = Flask(__name__)

UPLOAD_FOLDER = 'uploads'
SCREENSHOT_FOLDER = 'screenshots'
TEXT_FOLDER = 'text_content'

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

if not os.path.exists(SCREENSHOT_FOLDER):
    os.makedirs(SCREENSHOT_FOLDER)

if not os.path.exists(TEXT_FOLDER):
    os.makedirs(TEXT_FOLDER)

convertapi.api_secret = os.environ.get('ppt')

def read_pptx(ppt_file):
    T = ""
    # Function to extract text from slides
    def extract_text_from_slide(slide):
        text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text.append(run.text)
        return ' '.join(text)

    # Read the PowerPoint file
    prs = Presentation(ppt_file)

    # Dictionary to store text content by slide numbers
    slide_contents = {}

    # Extract text from each slide and write to a text file
    with open('textfile.txt', 'w') as f:
        for idx, slide in enumerate(prs.slides):
            slide_text = extract_text_from_slide(slide)
            slide_contents[idx] = slide_text
            f.write(slide_text + '\n')
            T += slide_text

    return T
    
def take_screenshot(file_path, screenshot_folder):
    print("taking screenshots")
    convertapi.api_secret = os.environ.get('ppt')
    convertapi.convert('jpg', {
        'File': file_path,
    }, from_format = 'ppt').save_files(screenshot_folder)


def check_images_presence(pptx_file):
    prs = Presentation(pptx_file)
    images_present = []

    for slide in prs.slides:
        image_found = False
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image_found = True
                break
        images_present.append(image_found)

    return images_present
from openai import OpenAI

def generate_topic_name(T, text):
    # Initialize the OpenAI client with your secure API key
    client = "OpenAI(api_key=os.environ.get('openai'))"

    # Create the completion request
    completion = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": f"give the name of the topic from this {T} these are my keywords {text}"}
        ]
    )

    # Extract and return the generated topic name
    return completion.choices[0].message.content

@app.after_request
async def add_cors_headers(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Methods', 'POST')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
    return response

@app.route('/upload', methods=['POST'])
async def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    if file and file.filename.endswith(('ppt', 'pptx')):
        filename = secure_filename(file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(file_path)

        # Take a screenshot of the uploaded PowerPoint file
        take_screenshot(file_path, SCREENSHOT_FOLDER)
        print("=========================ScreenShot complete============================================")
        # # # Check for images in the PowerPoint file
        images_presence = check_images_presence(file_path)
        print("=======================image presence  complete==========================================")
        print(images_presence)
        # # Extract text from the PowerPoint file
        print(file_path)
        T=read_pptx(file_path)
        print("=======================Text extraction Complete==========================================")
        
        Title=extract_slide_titles(file_path)
        print("=======================Title extraction complete==========================================")
        text=""
        for i in Title:
            text+=Title[i]+","
        topic_name=generate_topic_name(T,text)
        print("=======================Topic name==========================================")
        print(topic_name)
        explanations=generate_explanations(Title,topic_name)
        print(explanations)
        print("=======================Explainations complete==========================================")
        
        file_paths=sort_files_and_get_paths("screenshots")

        print("============================================================================")
        print(len(explanations))
        print(len(images_presence))
        print(len(Title))
        print(len(file_paths))
        
        explanations=generate_image_explanations(Title,images_presence,file_paths,model,explanations)
        print("======================image explaination================================")
        generate_and_save_audio_files(Title,explanations)
        print("======================Audio generated=====================================")
        return jsonify({'message': 'File uploaded successfully', 'filename': filename,
                        'images_present': images_presence}), 200
    else:
        return jsonify({'error': 'Please upload a PowerPoint file (PPT or PPTX)'}), 400

if __name__ == '__main__':
    app.run(debug=True,port=8000)